from flask import Flask, request, send_file, render_template
import os
from mistralai import Mistral
from pptx import Presentation

# Flask App
app = Flask(__name__)

# Initialize Mistral Client
api_key = os.getenv("MISTRAL_API_KEY")
client = Mistral(api_key=api_key)

@app.route("/")
def home():
    return render_template("ai-ppt.html")

@app.route("/generate", methods=["POST"])
def generate_ppt():
    topic = request.form["topic"]
    num_slides = int(11)

    # Generate slide content using Mistral
    slides_content = []
    for i in range(num_slides):
        response = client.agents.complete(
            agent_id="ag:ffcab88e:20241212:invfest:1a40a617",
            messages=[
                {
                    "role": "user",
                    "content": f"Buatkan Slide tentang: {topic}, slide {i + 1}",
                },
            ],
        )
        slide_text = response.choices[0].message.content
        slides_content.append(slide_text)

    # Create PPT
    ppt_file = create_ppt(topic, slides_content)

    return send_file(ppt_file, as_attachment=True)

def create_ppt(topic, slides_content):
    prs = Presentation()
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = topic
    slide.placeholders[1].text = "Generated by AI PPT Generator"

    # Content Slides
    for content in slides_content:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = topic
        slide.placeholders[1].text = content

    # Save PPT
    output_file = "generated_presentation.pptx"
    prs.save(output_file)
    return output_file

if __name__ == "__main__":
    app.run(debug=True)
